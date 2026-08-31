import os
import shutil
import logging
import tempfile
from pathlib import Path
from typing import Optional, List, Dict
from dataclasses import dataclass, field

logger = logging.getLogger('pptmaker.extractors.speaker')

MIN_PHOTO_HEIGHT_CM = 3.0
MIN_ASPECT_RATIO = 0.4
MAX_ASPECT_RATIO = 2.5


def _image_height_cm(img_path: str) -> float:
    """Return image physical height in cm. Uses DPI if available, else assumes 200 DPI."""
    try:
        from PIL import Image
        img = Image.open(img_path)
        dpi = img.info.get('dpi', None)
        h_px = img.height
        if dpi and dpi[1] and dpi[1] > 0:
            return h_px / dpi[1] * 2.54
        return h_px / 200 * 2.54
    except Exception:
        return 0.0


def _is_photo(img_path: str) -> bool:
    """Check if image looks like a photo (reasonable aspect ratio + min height)."""
    try:
        from PIL import Image
        img = Image.open(img_path)
        w, h = img.size
        ratio = w / h if h > w else h / w
        if ratio < MIN_ASPECT_RATIO or ratio > MAX_ASPECT_RATIO:
            return False
        return _image_height_cm(img_path) >= MIN_PHOTO_HEIGHT_CM
    except Exception:
        return False


_SESSION_PHOTO_DIR = None


def _get_session_photo_dir():
    global _SESSION_PHOTO_DIR
    if _SESSION_PHOTO_DIR is None or not os.path.exists(_SESSION_PHOTO_DIR):
        _SESSION_PHOTO_DIR = tempfile.mkdtemp(prefix='ppt_speaker_photos_')
    return _SESSION_PHOTO_DIR


def _crop_head_portrait(photo_path: str) -> bool:
    """照片裁剪为上半身像（头+肩+胸，脸约占画面高度 1/3）。返回是否已裁剪；任何失败保留原图。

    用 OpenCV Haar 级联检测人脸：检测不到人脸、或人脸已占画面 40% 以上
    （本来就是特写/证件照）时不裁剪；否则围绕人脸裁剪上半身区域。
    裁剪框一律限制在原图范围内——只裁剪，绝不扩图，也不放大分辨率。
    """
    import numpy as np
    import cv2
    from PIL import Image

    try:
        with open(photo_path, 'rb') as f:
            raw = f.read()
        # 用 imdecode 而非 imread：避免 Windows 下中文路径无法读取
        img = cv2.imdecode(np.frombuffer(raw, np.uint8), cv2.IMREAD_COLOR)
        if img is None:
            return False
        h, w = img.shape[:2]
        cascade = cv2.CascadeClassifier(
            cv2.data.haarcascades + 'haarcascade_frontalface_default.xml')
        gray = cv2.equalizeHist(cv2.cvtColor(img, cv2.COLOR_BGR2GRAY))

        # 先严格、后宽松两轮检测，尽量找出人脸
        faces = cascade.detectMultiScale(
            gray, scaleFactor=1.1, minNeighbors=5,
            minSize=(int(min(w, h) * 0.06), int(min(w, h) * 0.06)))
        if len(faces) == 0:
            faces = cascade.detectMultiScale(
                gray, scaleFactor=1.05, minNeighbors=3,
                minSize=(int(min(w, h) * 0.04), int(min(w, h) * 0.04)))
        if len(faces) == 0:
            return False
        # 取面积最大的人脸（通常是照片主体）
        fx, fy, fw, fh = max(faces, key=lambda f: f[2] * f[3])

        # 人脸已占画面高度 40% 以上 → 已是特写，不裁剪
        if fh >= 0.4 * h:
            return False

        # 上半身区域：以人脸为中心，上留头发、下至胸肩、两侧含肩
        # （高度 3.0×人脸高 → 脸约占 1/3；宽度 2.4×人脸宽 ≈ 1.4×头宽，含肩）
        cx = fx + fw / 2.0
        x0 = int(cx - 1.2 * fw)
        x1 = int(cx + 1.2 * fw)
        y0 = int(fy - 0.6 * fh)
        y1 = int(fy + 2.4 * fh)
        # 限制在原图范围内——只裁剪，绝不扩图
        x0 = max(0, x0)
        y0 = max(0, y0)
        x1 = min(w, x1)
        y1 = min(h, y1)
        if x1 <= x0 or y1 <= y0:
            return False

        with Image.open(photo_path) as im:
            im.load()
            fmt = im.format or 'JPEG'
            cropped = im.crop((x0, y0, x1, y1))
        # 按原格式、原始分辨率保存，不放大；JPEG 用高质量减少重编码损失
        save_kwargs = {'quality': 95, 'subsampling': 0} if fmt == 'JPEG' else {}
        cropped.save(photo_path, format=fmt, **save_kwargs)
        logger.info('照片已裁剪为上半身: %dx%d → (%d,%d,%d,%d)，人脸框 (%d,%d,%d,%d)',
                    w, h, x0, y0, x1, y1, fx, fy, fw, fh)
        return True
    except Exception as e:
        logger.warning('头像裁剪失败，保留原图: %s', e)
        return False


@dataclass
class Speaker:
    name: str
    photo_path: str = ''
    bio: str = ''
    bio_en: str = ''
    institution: str = ''
    title: str = ''
    photo_original_path: str = ''  # 裁剪前的完整照片（未裁剪时为空，等于 photo_path）


def _read_zip_entry_raw(zf, name):
    """Read a ZIP entry even if its CRC is corrupted.

    Returns the uncompressed data, or ``None`` if unrecoverable.
    """
    import struct
    import zlib

    info = zf.getinfo(name)
    fp = zf.fp
    fp.seek(info.header_offset)
    local_header = fp.read(30)
    if len(local_header) < 30 or local_header[:4] != b'PK\x03\x04':
        return None
    filename_len = struct.unpack('<H', local_header[26:28])[0]
    extra_len = struct.unpack('<H', local_header[28:30])[0]
    data_offset = info.header_offset + 30 + filename_len + extra_len
    fp.seek(data_offset)
    compressed = fp.read(info.compress_size)

    if info.compress_type == 0:   # stored (no compression)
        return compressed
    if info.compress_type == 8:   # deflated
        try:
            return zlib.decompress(compressed, -15)
        except Exception:
            return None
    return None


def _extract_images_from_docx(file_path: str, output_dir: str) -> List[str]:
    """Extract embedded images from a DOCX file. Returns list of image paths."""
    from docx import Document

    try:
        doc = Document(file_path)
    except Exception:
        return _extract_images_from_docx_zip_fallback(file_path, output_dir)

    paths = []
    for i, rel in enumerate(doc.part.rels.values()):
        if 'image' in rel.reltype:
            ext = rel.target_ref.split('.')[-1].split('?')[0]
            if ext not in ('jpeg', 'jpg', 'png', 'gif', 'bmp'):
                ext = 'png'
            img_path = os.path.join(output_dir, f'docx_img_{i}.{ext}')
            try:
                with open(img_path, 'wb') as f:
                    f.write(rel.target_part.blob)
                paths.append(img_path)
            except Exception:
                continue
    return paths


def _extract_images_from_docx_zip_fallback(file_path: str, output_dir: str) -> List[str]:
    """Fallback: extract images from DOCX ZIP directly, including entries
    with corrupted CRC (raw byte extraction)."""
    import zipfile
    import struct
    import zlib

    paths = []
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            for i, name in enumerate(zf.namelist()):
                if not name.startswith('word/media/'):
                    continue
                try:
                    data = zf.read(name)
                except Exception:
                    # CRC error or other read failure —
                    # try to read the raw bytes bypassing the CRC check
                    data = _read_zip_entry_raw(zf, name)
                    if data is None:
                        continue
                ext = name.rsplit('.', 1)[-1].lower()
                if ext not in ('jpeg', 'jpg', 'png', 'gif', 'bmp'):
                    ext = 'png'
                img_path = os.path.join(output_dir, f'docx_img_{i}.{ext}')
                with open(img_path, 'wb') as f:
                    f.write(data)
                paths.append(img_path)
    except Exception:
        pass
    return paths


def _extract_images_from_pptx(file_path: str, output_dir: str) -> List[str]:
    """Extract embedded images from a PPTX file. Returns list of image paths.

    Covers both direct PICTURE shapes and placeholders/autoshapes that
    contain picture fills (common when a user inserts a photo into a
    content placeholder).
    """
    from pptx import Presentation
    from pptx.oxml.ns import qn

    prs = Presentation(file_path)
    paths = []
    seen = set()  # deduplicate by rId
    nsm = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            img_blob = None
            img_ct = None

            # 1) Direct PICTURE shape
            if shape.shape_type == 13 and shape.image:
                img_blob = shape.image.blob
                img_ct = shape.image.content_type

            # 2) Any shape with a picture fill (blipFill) —
            #    covers placeholders, autoshapes, etc.
            if img_blob is None:
                blips = shape._element.findall('.//a:blip', nsm)
                for blip in blips:
                    embed = blip.get(qn('r:embed'))
                    if not embed or embed in seen:
                        continue
                    if embed in slide.part.rels:
                        rel = slide.part.rels[embed]
                        img_part = getattr(rel, 'target_part', rel)
                        blob = getattr(img_part, 'blob', None) or getattr(img_part, '_blob', None)
                        if blob:
                            img_blob = blob
                            img_ct = getattr(img_part, 'content_type', 'image/png')
                            seen.add(embed)
                            break

            if img_blob is None:
                continue

            ext = (img_ct or 'image/png').split('/')[-1]
            if ext == 'jpeg':
                ext = 'jpg'
            if ext not in ('jpg', 'jpeg', 'png', 'gif', 'bmp'):
                ext = 'png'
            img_path = os.path.join(output_dir, f'pptx_img_{len(paths)}.{ext}')
            with open(img_path, 'wb') as f:
                f.write(img_blob)
            paths.append(img_path)
    return paths


def _extract_images_from_pdf(file_path: str, output_dir: str) -> List[str]:
    """Extract embedded images from a PDF file via PyPDF2 page XObjects."""
    from PyPDF2 import PdfReader

    paths = []
    try:
        reader = PdfReader(file_path)
    except Exception:
        return paths

    for page_idx, page in enumerate(reader.pages):
        resources = page.get('/Resources', {})
        if not resources:
            continue
        xobjects = resources.get('/XObject', {})
        if not xobjects:
            continue
        for img_idx, (name, obj) in enumerate(xobjects.items()):
            try:
                xobj = obj.get_object()
            except Exception:
                continue
            if xobj.get('/Subtype', '') != '/Image':
                continue
            try:
                data = xobj.get_data()
            except Exception:
                continue
            if not data:
                continue
            # Determine extension from the image filter
            f = xobj.get('/Filter', '')
            if isinstance(f, list):
                f = f[0] if f else ''
            f_name = f.get('/Name', '') if hasattr(f, 'get') else str(f)
            if 'DCTDecode' in f_name or 'DCT' in f_name:
                ext = 'jpg'
            elif 'JPXDecode' in f_name:
                ext = 'jp2'
            else:
                ext = 'png'
            if ext not in ('jpg', 'jpeg', 'png', 'gif', 'bmp'):
                ext = 'png'
            img_path = os.path.join(
                output_dir, f'pdf_p{page_idx}_i{img_idx}.{ext}')
            with open(img_path, 'wb') as fh:
                fh.write(data)
            paths.append(img_path)
    return paths


def _extract_text_from_docx(file_path: str) -> str:
    from docx import Document
    try:
        doc = Document(file_path)
    except Exception:
        return _extract_text_from_docx_zip_fallback(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return '\n'.join(paragraphs)


def _extract_text_from_docx_zip_fallback(file_path: str) -> str:
    """Fallback: extract text from word/document.xml directly, tolerating corrupt media."""
    import zipfile
    from xml.etree import ElementTree
    import re

    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            xml_bytes = zf.read('word/document.xml')
    except Exception:
        raise RuntimeError(f'无法读取DOCX文件: {file_path}')

    root = ElementTree.fromstring(xml_bytes)
    ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
    texts = []
    for t in root.iter('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t'):
        if t.text:
            texts.append(t.text)
    # Join with newlines at paragraph breaks
    full = ''.join(texts)
    # Split on paragraph markers
    paragraphs = [p.strip() for p in re.split(r'\n\s*', full) if p.strip()]
    return '\n'.join(paragraphs)


def _extract_text_from_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
    return '\n'.join(texts)


def _extract_text_from_pdf(file_path: str) -> str:
    from PyPDF2 import PdfReader
    reader = PdfReader(file_path)
    texts = []
    for page in reader.pages:
        t = page.extract_text()
        if t:
            texts.append(t)
    return '\n'.join(texts)


def _extract_text_from_doc_ole(file_path: str) -> str:
    """Extract text from binary .doc (OLE compound) file."""
    import olefile
    import re

    try:
        ole = olefile.OleFileIO(file_path)
    except Exception:
        return ''

    parts = []
    try:
        if ole.exists('WordDocument'):
            data = ole.openstream('WordDocument').read()
            # Decode as UTF-16LE (most .doc files use Unicode storage)
            decoded = data.decode('utf-16-le', errors='ignore')
            # Remove non-text characters; keep only Chinese, ASCII, punctuation
            cleaned = re.sub(
                r'[^一-鿿　-〿＀-￯'
                r'a-zA-Z0-9\s\.\,\;\:\!\?\(\)\[\]\{\}\-\+\/\@'
                r'\#\$\%\&\*\"\'\<\>\=\~\`\_\|\n\r'
                r'‐-⁯℀-⅏]+',
                '\n', decoded)
            # Collapse multiple newlines
            cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
            # Only keep lines that contain at least 2 consecutive CJK chars
            # (filters out binary noise that decoded to random symbols)
            filtered = []
            for line in cleaned.split('\n'):
                line = line.strip()
                if not line:
                    continue
                # Keep lines with at least 2 consecutive Chinese chars
                # or meaningful ASCII text (3+ words)
                if re.search(r'[一-鿿]{2,}', line) or len(line.split()) >= 3:
                    filtered.append(line)
            parts.append('\n'.join(filtered))
    finally:
        ole.close()

    return '\n'.join(p.strip() for p in '\n'.join(parts).split('\n') if p.strip())


def _extract_images_from_doc_ole(file_path: str, output_dir: str) -> List[str]:
    """Extract images from binary .doc (OLE compound) file by scanning
    all streams for JPEG/PNG signatures."""
    import olefile

    paths = []
    # Image signatures
    sigs = {
        b'\xff\xd8\xff': 'jpg',    # JPEG
        b'\x89PNG\r\n\x1a\n': 'png',  # PNG
        b'GIF8': 'gif',           # GIF
        b'BM': 'bmp',             # BMP
    }

    try:
        ole = olefile.OleFileIO(file_path)
    except Exception:
        return paths

    try:
        for stream_name in ole.listdir():
            sname = '/'.join(stream_name)
            try:
                data = ole.openstream(sname).read()
            except Exception:
                continue
            # Search for image signatures in stream data
            for sig_bytes, ext in sigs.items():
                start = 0
                while True:
                    idx = data.find(sig_bytes, start)
                    if idx == -1:
                        break
                    # Try to read the image
                    if ext == 'jpg':
                        # JPEG: find end marker FF D9
                        end = data.find(b'\xff\xd9', idx + 2)
                        if end == -1:
                            end = len(data)
                        else:
                            end += 2
                    elif ext == 'png':
                        # PNG: find IEND marker
                        end = data.find(b'IEND\xae\x42\x60\x82', idx)
                        if end == -1:
                            end = len(data)
                        else:
                            end += 8
                    elif ext == 'gif':
                        end = data.find(b'\x00\x3b', idx)  # GIF terminator
                        if end == -1:
                            end = len(data)
                        else:
                            end += 2
                    else:
                        end = min(idx + 10 * 1024 * 1024, len(data))  # cap at 10MB

                    img_data = data[idx:end]
                    if len(img_data) > 256:  # skip tiny fragments
                        img_path = os.path.join(
                            output_dir, f'doc_img_{len(paths)}.{ext}')
                        with open(img_path, 'wb') as f:
                            f.write(img_data)
                        paths.append(img_path)
                    start = idx + 1
    finally:
        ole.close()

    return paths


def _extract_text_from_ppt_ole(file_path: str) -> str:
    """Extract text from binary .ppt (OLE compound) file.

    Searches for TextCharsAtom (0x0FA0) and TextBytesAtom (0x0FA8)
    records within the PowerPoint Document stream and extracts their
    text content.
    """
    import olefile
    import struct
    import re

    try:
        ole = olefile.OleFileIO(file_path)
    except Exception:
        return ''

    texts = []
    try:
        if not ole.exists('PowerPoint Document'):
            return ''
        data = ole.openstream('PowerPoint Document').read()

        # Search for TextCharsAtom (recType=0x0FA0) and
        # TextBytesAtom (recType=0x0FA8) record headers.
        pos = 0
        while pos < len(data) - 8:
            rec_type = struct.unpack_from('<H', data, pos + 2)[0]

            if rec_type in (0x0FA0, 0x0FA8):
                rec_len = struct.unpack_from('<I', data, pos + 4)[0]
                text_start = pos + 8
                if (text_start + rec_len <= len(data)
                        and 0 < rec_len < 50000):
                    if rec_type == 0x0FA0:  # TextCharsAtom — UTF-16LE
                        try:
                            t = data[text_start:text_start + rec_len].decode('utf-16-le', errors='ignore')
                            t = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', t)
                            if t.strip():
                                texts.append(t.strip())
                        except Exception:
                            pass
                    else:  # TextBytesAtom — single-byte encoded
                        try:
                            t = data[text_start:text_start + rec_len].decode('cp1252', errors='ignore')
                            if t.strip():
                                texts.append(t.strip())
                        except Exception:
                            pass
                    pos = text_start + rec_len
                else:
                    pos += 1
            else:
                pos += 1
    finally:
        ole.close()

    # Deduplicate and join
    seen = set()
    result = []
    for t in texts:
        if t not in seen:
            seen.add(t)
            result.append(t)
    return '\n'.join(result)


def _extract_images_from_ppt_ole(file_path: str, output_dir: str) -> List[str]:
    """Extract images from binary .ppt (OLE compound) file.
    PPT stores images in the \"Pictures\" stream."""
    import olefile
    import struct

    paths = []
    try:
        ole = olefile.OleFileIO(file_path)
    except Exception:
        return paths

    try:
        if ole.exists('Pictures'):
            data = ole.openstream('Pictures').read()
            # .ppt Pictures stream contains embedded images with headers
            # Scan for JPEG/PNG signatures
            sigs = {
                b'\xff\xd8\xff': ('jpg', b'\xff\xd9'),
                b'\x89PNG\r\n\x1a\n': ('png', b'IEND\xae\x42\x60\x82'),
            }
            for sig_bytes, (ext, end_marker) in sigs.items():
                start = 0
                while True:
                    idx = data.find(sig_bytes, start)
                    if idx == -1:
                        break
                    end = data.find(end_marker, idx + 2)
                    if end == -1:
                        end = len(data)
                    else:
                        end += len(end_marker)
                    img_data = data[idx:end]
                    if len(img_data) > 256:
                        img_path = os.path.join(
                            output_dir, f'ppt_img_{len(paths)}.{ext}')
                        with open(img_path, 'wb') as f:
                            f.write(img_data)
                        paths.append(img_path)
                    start = idx + 1
    finally:
        ole.close()

    return paths


def _resolve_institution(short_name: str, api_key: str) -> str:
    """Use DeepSeek to expand a hospital abbreviation to its full name.

    Returns the full name (ending with 医院), or empty string if
    resolution fails.
    """
    from utils.deepseek_client import structure_text

    prompt = (
        '请将以下医院简称扩展为完整官方全称，必须以"医院"结尾。\n'
        '如果这是某个知名医院的简称或俗称，请返回其完整名称。\n'
        '如果无法确定全称，只返回"UNKNOWN"。\n'
        '只返回全称文字，不要解释、不要加引号。'
    )
    try:
        result = structure_text(short_name, prompt, api_key)
        result = result.strip().strip('"\'').strip()
        if result and result != 'UNKNOWN' and result.endswith('医院'):
            return result
    except Exception:
        pass
    return ''


def _translate_bio_en(bio: str, api_key: str) -> str:
    """Translate a Chinese bio into English via DeepSeek.

    Returns the translated multi-line text, or empty string on failure.
    """
    import json
    from utils.deepseek_client import structure_text

    prompt = (
        '请将以下中文履历逐行翻译成英文，返回JSON格式：'
        '{"bio_en": ["第一行英文", "第二行英文"]}\n'
        '保持原文的行结构：一行对应一行，不要合并或拆分。\n'
        '医学专业术语使用标准英文表达。\n'
        '不要添加任何原文没有的内容。\n\n'
        '履历：\n' + bio
    )
    try:
        result = structure_text(bio, prompt, api_key)
        data = json.loads(result)
        lines = data.get('bio_en', [])
        if isinstance(lines, str):
            return lines
        return '\n'.join(str(l) for l in lines)
    except Exception:
        return ''


def extract_speaker(data, api_key: str, filename: Optional[str] = None) -> Speaker:
    """
    Extract speaker info from a file (docx/doc/pptx/ppt/pdf).
    Uses python libraries for text/image extraction, then DeepSeek to structure.
    Binary .doc/.ppt use OLE-based extraction.

    data: file bytes (web upload) or a local path (desktop).
    filename: original filename, used for extension detection and name
              inference. Required for bytes input; derived from the path
              when omitted.
    """
    if filename is None:
        if isinstance(data, (str, os.PathLike)):
            filename = str(data)
        else:
            raise ValueError('bytes 输入必须提供 filename 参数')

    # Materialize bytes to a temp file with the right suffix; desktop paths
    # pass through unchanged. Internal helpers keep working on real paths.
    _tmp_path = None
    if isinstance(data, (bytes, bytearray)):
        _tmp_path = tempfile.NamedTemporaryFile(
            suffix=Path(filename).suffix.lower(), delete=False).name
        with open(_tmp_path, 'wb') as f:
            f.write(bytes(data))
        file_path = _tmp_path
    else:
        if not isinstance(data, (str, os.PathLike)):
            raise TypeError(f'不支持的文件输入类型: {type(data).__name__}（应为 bytes 或本地路径）')
        file_path = str(data)

    try:
        size_hint = len(data) if isinstance(data, (bytes, bytearray)) else os.path.getsize(data)
    except OSError:
        size_hint = -1
    logger.info('开始提取演讲者: %s (%d bytes)', filename, size_hint)

    ext = Path(file_path).suffix.lower()
    work_dir = tempfile.mkdtemp(prefix='spk_')
    import re
    stem = Path(filename).stem
    # Pattern: optional prefix (number/underscore) + Chinese name + suffix
    name_from_file = re.sub(r'^[\d_]+', '', stem)
    name_from_file = re.sub(r'(简介|简历|介绍|资料|个人|履历)$', '', name_from_file)
    if not name_from_file:
        name_from_file = stem

    try:
        # Extract text — try XML-based first, fall back to OLE for .doc/.ppt
        if ext == '.docx':
            raw_text = _extract_text_from_docx(file_path)
            img_paths = _extract_images_from_docx(file_path, work_dir)
        elif ext == '.doc':
            # Try python-docx first (some .doc files are actually .docx);
            # binary .doc files will raise an exception, falling back to OLE.
            try:
                raw_text = _extract_text_from_docx(file_path)
            except Exception:
                raw_text = ''
            if not raw_text.strip():
                raw_text = _extract_text_from_doc_ole(file_path)
            try:
                img_paths = _extract_images_from_docx(file_path, work_dir)
            except Exception:
                img_paths = []
            if not img_paths:
                img_paths = _extract_images_from_doc_ole(file_path, work_dir)
        elif ext == '.pptx':
            raw_text = _extract_text_from_pptx(file_path)
            img_paths = _extract_images_from_pptx(file_path, work_dir)
        elif ext == '.ppt':
            # Try python-pptx first (some .ppt files are actually .pptx);
            # binary .ppt files will raise an exception, falling back to OLE.
            try:
                raw_text = _extract_text_from_pptx(file_path)
            except Exception:
                raw_text = ''
            if not raw_text.strip():
                raw_text = _extract_text_from_ppt_ole(file_path)
            try:
                img_paths = _extract_images_from_pptx(file_path, work_dir)
            except Exception:
                img_paths = []
            if not img_paths:
                img_paths = _extract_images_from_ppt_ole(file_path, work_dir)
        elif ext == '.pdf':
            raw_text = _extract_text_from_pdf(file_path)
            img_paths = _extract_images_from_pdf(file_path, work_dir)
        else:
            raise ValueError(f'不支持的文件格式: {ext}')

        if not raw_text:
            raise RuntimeError(f'未能从文件中提取到文字: {file_path}')

        # Drop the resume header line (e.g. "朱颜铂 教授") so the name/title
        # doesn't pollute the bio
        # PPTX paragraph soft-breaks use \x0b — treat as line separators too
        lines = re.split(r'[\n\x0b]', raw_text)
        first = lines[0].strip() if lines else ''
        if first and name_from_file:
            if first == name_from_file:
                del lines[0]
            else:
                rest = re.sub(rf'^{re.escape(name_from_file)}\s*', '', first)
                # Name + short title (教授/主任医师/院长 etc.)
                if rest != first and len(rest) <= 8:
                    del lines[0]
        raw_text = '\n'.join(lines)

        # Use DeepSeek to structure the text
        from utils.deepseek_client import structure_text

        prompt = (
            '请从以下文本中识别演讲者的信息，返回JSON格式：\n'
            '{"name": "姓名", "institution": "所属医院", '
            '"bio": "履历列表"}\n'
            '\n'
            f'提示：文件名为「{stem}」，该演讲者的姓名很可能是「{name_from_file}」。\n'
            '如果文本中找不到明确的姓名，请使用提示中的姓名。\n'
            'name字段必须返回非空姓名：中文姓名返回2-3个汉字，外籍讲者返回英文全名'
            '（如"Theo M. de Reijke"），不要包含职务和头衔。\n'
            'institution字段从文本中提取医院名称，如"重庆大学附属肿瘤医院"（只保留一个最完整的医院名称）。\n'
            'bio字段返回JSON字符串数组，每项为一行履历。\n'
            '\n'
            '核心原则：尽量保持原文的行结构和自然分段，不要随意按顿号、逗号拆分成多行。\n'
            '原文中同一行/同一段的内容保持在同一行，如原文将所有学会任职和编委写在一行，\n'
            '就全部保留在一行（如"担任中国医师协会...会长、北京医学会...委员，UroPrecison主编、\n'
            'J of Urology编委、《中华医学杂志》副总编"）。\n'
            '只在原文有明确换行或分段时才可以分成多行。\n'
            '\n'
            '数组内容按以下顺序整理（保留全部信息，不省略）：\n'
            '1. 领导职务——如有多项用顿号连接在一行\n'
            '2. 学位职称——学位、职称、导师资格合并一行\n'
            '3. 临床专长——临床和研究方向描述\n'
            '4. 科研成果——承担课题、发表论文、获奖等，各部分自然分段\n'
            '5. 荣誉称号——每个单独一行\n'
            '6. 学会任职及期刊编委——学会任职和期刊编委职务，原文在同一段的保持在一行\n'
            '删除"共"、"荣获"等多余连接词，保留核心内容。\n'
            '如果文本中包含多个人的信息，只提取主要人物。'
        )
        result = structure_text(raw_text, prompt, api_key)
        import json
        data = json.loads(result)

        # 最终姓名：DeepSeek 返回空/null/缺失时回退到文件名推断，
        # 保证解析结果姓名框不为空（照片文件名、bio_en 判断等均用同一姓名）
        name = (data.get('name') or '').strip()
        if not name:
            name = name_from_file
            logger.warning('DeepSeek 未返回姓名，回退为文件名推断: %s（文件: %s）',
                           name, filename)

        # Pick the largest qualifying photo (height >= 3cm), copy to session temp dir
        photo_path = ''
        photo_original_path = ''
        qualifying = [(p, os.path.getsize(p)) for p in img_paths if _is_photo(p)]
        if qualifying:
            best = max(qualifying, key=lambda x: x[1])[0]
            ext = os.path.splitext(best)[1]
            safe_name = name.replace('/', '_').replace('\\', '_')
            dest = os.path.join(_get_session_photo_dir(), f'{safe_name}{ext}')
            shutil.copy2(best, dest)
            photo_path = dest
            # 照片裁剪为上半身像（人脸过小才裁；失败保留原图，不影响提取流程）
            try:
                if _crop_head_portrait(dest):
                    # 裁剪成功时另存一份完整原图，供「查看完整照片」
                    photo_original_path = os.path.join(
                        _get_session_photo_dir(), f'{safe_name}_完整照片{ext}')
                    shutil.copy2(best, photo_original_path)
            except Exception:
                pass

        # Resolve institution full name if it doesn't end with 医院
        institution = data.get('institution', '')
        if institution and not institution.endswith('医院'):
            try:
                resolved = _resolve_institution(institution, api_key)
                if resolved:
                    institution = resolved
            except Exception:
                pass  # Keep original if resolution fails

        bio = data.get('bio', '')
        if isinstance(bio, list):
            bio = '\n'.join(bio)

        # English-named speakers get a translated bio for the second page
        bio_en = ''
        if bio and re.search(r'[A-Za-z]', name):
            try:
                bio_en = _translate_bio_en(bio, api_key)
            except Exception:
                pass

        speaker = Speaker(
            name=name,
            photo_path=photo_path,
            bio=bio,
            bio_en=bio_en,
            institution=institution,
            title='教授',
            photo_original_path=photo_original_path,
        )
        logger.info('演讲者提取完成: %s → %s', filename, speaker.name)
        return speaker
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)
        if _tmp_path:
            try:
                os.unlink(_tmp_path)
            except OSError:
                pass
