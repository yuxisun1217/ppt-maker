import json
import os
import base64
import urllib.request
import urllib.error
import urllib.parse
from typing import List
from dataclasses import dataclass


OCR_SPACE_URL = 'https://api.ocr.space/parse/image'


@dataclass
class AgendaItem:
    order: int
    time_slot: str
    session_title_cn: str
    session_title_en: str
    speaker_name: str
    host: str          # 主持人/主持嘉宾
    institution: str
    item_type: str     # 'speech' | 'panel' | 'opening' | 'closing'


AGENDA_SCHEMA = (
    '返回JSON数组，每个环节包含：\n'
    '  "order": 顺序号(从1开始，section标题也要编入序号),\n'
    '  "time_slot": "时间段",\n'
    '  "session_title_cn": "中文标题",\n'
    '  "session_title_en": "英文标题(保留原文)",\n'
    '  "speaker_name": "演讲者姓名",\n'
    '  "host": "主持人姓名(panel讨论的主持嘉宾，普通演讲留空)",\n'
    '  "institution": "机构",\n'
    '  "item_type": "类型: section/opening/speech/panel/closing"\n'
    '\n'
    '重要：日程中的章节标题（如"泌见前沿，云端共话""渝见前沿，加备精彩"等大字分隔标题）'
    '也必须提取，type设为"section"，只填order和session_title_cn，其他字段留空。\n'
    '\n'
    '返回格式：{"agenda": [...]}\n'
    '注意：严格按时间顺序排列。\n'
    'speaker_name和host只保留姓名，删除所有职称/职务（如"教授""主任""院长"等），只留纯姓名。\n'
)

OCR_PROMPT = (
    '以下文字是从会议日程图片中通过OCR识别提取的。'
    '部分中文字符可能识别不准确，但时间、英文单词和结构信息基本可靠。\n'
    '请根据这些文字还原会议日程中的全部内容。\n'
    + AGENDA_SCHEMA +
    '\nOCR文字：\n'
)

DOC_PROMPT = (
    '以下文字是从会议日程文件中直接提取的，文字准确度较高。\n'
    '请根据这些文字识别会议日程中的全部内容。\n'
    + AGENDA_SCHEMA +
    '\n日程文字：\n'
)


def _compress_image(image_path: str, max_bytes: int = 200 * 1024) -> str:
    """Compress image to under max_bytes. Returns path to compressed temp file."""
    from PIL import Image
    import tempfile

    # Check original size
    orig_size = os.path.getsize(image_path)
    if orig_size <= max_bytes:
        return image_path

    img = Image.open(image_path)
    if img.mode in ('RGBA', 'P', 'LA'):
        img = img.convert('RGB')

    # Try JPEG compression with decreasing quality
    out = os.path.join(tempfile.mkdtemp(prefix='ocr_compressed_'), 'compressed.jpg')
    for quality in (85, 70, 55, 40, 25):
        img.save(out, 'JPEG', quality=quality, optimize=True)
        if os.path.getsize(out) <= max_bytes:
            return out

    # Still too big — scale down dimensions by 20% each step
    w, h = img.size
    for _ in range(6):
        w, h = int(w * 0.8), int(h * 0.8)
        resized = img.resize((w, h), Image.LANCZOS)
        for quality in (70, 50, 30):
            resized.save(out, 'JPEG', quality=quality, optimize=True)
            if os.path.getsize(out) <= max_bytes:
                return out

    return out  # best effort


def _ocr_image(image_path: str, ocr_api_key: str) -> str:
    """Extract text from image using OCR.space API."""
    compressed_path = _compress_image(image_path)
    with open(compressed_path, 'rb') as f:
        img_b64 = base64.b64encode(f.read()).decode('utf-8')

    payload = urllib.parse.urlencode({
        'base64Image': f'data:image/jpeg;base64,{img_b64}',
        'language': 'chs',
        'isTable': 'true',
        'OCREngine': '3',
    }).encode('utf-8')

    req = urllib.request.Request(OCR_SPACE_URL, data=payload, method='POST')
    req.add_header('apikey', ocr_api_key)
    req.add_header('Content-Type', 'application/x-www-form-urlencoded')

    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            data = json.loads(resp.read().decode('utf-8'))
    except urllib.error.HTTPError as e:
        body = e.read().decode('utf-8', errors='ignore')
        raise RuntimeError(f'OCR.space API 错误 ({e.code}): {body}') from e

    if data.get('IsErroredOnProcessing', True):
        err_msg = data.get('ErrorMessage', ['Unknown error'])
        raise RuntimeError(f'OCR.space 处理失败: {err_msg}')

    results = data.get('ParsedResults', [])
    if not results:
        raise RuntimeError('OCR.space 未返回识别结果')

    text = results[0].get('ParsedText', '')
    if not text.strip():
        raise RuntimeError('OCR.space 返回空文本')
    return text


def _extract_text_from_docx(file_path: str) -> str:
    """Extract all text from a DOCX file."""
    from docx import Document
    doc = Document(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    # Also extract table content
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(' | '.join(cells))
    return '\n'.join(paragraphs)


def _extract_text_from_pptx(file_path: str) -> str:
    """Extract all text from a PPTX file."""
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        texts.append(' | '.join(cells))
    return '\n'.join(texts)


def _extract_text_from_xlsx(file_path: str) -> str:
    """Extract all text from an XLSX file, preserving table structure."""
    from openpyxl import load_workbook
    wb = load_workbook(file_path, data_only=True)
    all_lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        # If multiple sheets, add a header
        if len(wb.sheetnames) > 1:
            all_lines.append(f'--- Sheet: {sheet_name} ---')
        for row in ws.iter_rows(values_only=True):
            # Join non-empty cell values with separator
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                all_lines.append('\t'.join(cells))
    wb.close()
    return '\n'.join(all_lines)


def _extract_text_from_doc_ole(file_path: str) -> str:
    """Extract text from binary .doc (OLE) file — agenda context."""
    import olefile
    import re

    try:
        ole = olefile.OleFileIO(file_path)
    except Exception:
        return ''

    parts = []
    try:
        if ole.exists('WordDocument'):
            data = ole.openstream('WordDocument').read()
            decoded = data.decode('utf-16-le', errors='ignore')
            cleaned = re.sub(
                r'[^一-鿿　-〿＀-￯'
                r'a-zA-Z0-9\s\.\,\;\:\!\?\(\)\[\]\{\}\-\+\/\@'
                r'\#\$\%\&\*\"\'\<\>\=\~\`\_\|\n\r'
                r'‐-⁯℀-⅏]+',
                '\n', decoded)
            cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
            # Only keep lines with >=2 CJK chars or >=3 ASCII words
            filtered = []
            for line in cleaned.split('\n'):
                line = line.strip()
                if not line:
                    continue
                if re.search(r'[一-鿿]{2,}', line) or len(line.split()) >= 3:
                    filtered.append(line)
            parts.append('\n'.join(filtered))
    finally:
        ole.close()

    return '\n'.join(p.strip() for p in '\n'.join(parts).split('\n') if p.strip())


def _extract_text_from_ppt_ole(file_path: str) -> str:
    """Extract text from binary .ppt (OLE) file — agenda context.

    Searches for TextCharsAtom (0x0FA0) and TextBytesAtom (0x0FA8)
    records within the PowerPoint Document stream.
    """
    import olefile
    import struct
    import re

    try:
        ole = olefile.OleFileIO(file_path)
    except Exception:
        return ''

    texts = []
    try:
        if not ole.exists('PowerPoint Document'):
            return ''
        data = ole.openstream('PowerPoint Document').read()

        pos = 0
        while pos < len(data) - 8:
            rec_type = struct.unpack_from('<H', data, pos + 2)[0]

            if rec_type in (0x0FA0, 0x0FA8):
                rec_len = struct.unpack_from('<I', data, pos + 4)[0]
                text_start = pos + 8
                if (text_start + rec_len <= len(data)
                        and 0 < rec_len < 50000):
                    if rec_type == 0x0FA0:  # UTF-16LE
                        try:
                            t = data[text_start:text_start + rec_len].decode('utf-16-le', errors='ignore')
                            t = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', t)
                            if t.strip():
                                texts.append(t.strip())
                        except Exception:
                            pass
                    else:  # single-byte
                        try:
                            t = data[text_start:text_start + rec_len].decode('cp1252', errors='ignore')
                            if t.strip():
                                texts.append(t.strip())
                        except Exception:
                            pass
                    pos = text_start + rec_len
                else:
                    pos += 1
            else:
                pos += 1
    finally:
        ole.close()

    seen = set()
    result = []
    for t in texts:
        if t not in seen:
            seen.add(t)
            result.append(t)
    return '\n'.join(result)


def _parse_agenda_json(result: str) -> List[AgendaItem]:
    """Parse DeepSeek JSON response into AgendaItem list."""
    data = json.loads(result)
    items = data.get('agenda', [])
    agenda = []
    for item in items:
        agenda.append(AgendaItem(
            order=int(item.get('order', 0)),
            time_slot=str(item.get('time_slot', '')),
            session_title_cn=str(item.get('session_title_cn', '')),
            session_title_en=str(item.get('session_title_en', '')),
            speaker_name=str(item.get('speaker_name', '')),
            host=str(item.get('host', '')),
            institution=str(item.get('institution', '')),
            item_type=str(item.get('item_type', 'speech')),
        ))
    return agenda


def extract_agenda(file_path: str, api_key: str, ocr_api_key: str = '') -> List[AgendaItem]:
    """
    Extract agenda from image, PPTX/PPT, DOCX/DOC, or XLSX file.

    - Image (jpg/png): compress → OCR.space → DeepSeek
    - PPTX/PPT/DOCX/DOC/XLSX: extract text directly → DeepSeek (no OCR needed)
    - Binary .doc/.ppt fall back to OLE-based extraction
    """
    from utils.deepseek_client import structure_text

    ext = os.path.splitext(file_path)[1].lower()

    if ext in ('.jpg', '.jpeg', '.png'):
        if not ocr_api_key:
            raise RuntimeError('未配置 OCR.space API Key，请先在注册时填写或在账户设置中更新')
        raw_text = _ocr_image(file_path, ocr_api_key)
        result = structure_text(raw_text, OCR_PROMPT, api_key)
    elif ext in ('.pptx', '.ppt', '.docx', '.doc', '.xlsx'):
        if ext == '.pptx':
            raw_text = _extract_text_from_pptx(file_path)
        elif ext == '.ppt':
            try:
                raw_text = _extract_text_from_pptx(file_path)
            except Exception:
                raw_text = ''
            if not raw_text.strip():
                raw_text = _extract_text_from_ppt_ole(file_path)
        elif ext == '.xlsx':
            raw_text = _extract_text_from_xlsx(file_path)
        elif ext == '.docx':
            raw_text = _extract_text_from_docx(file_path)
        else:  # .doc
            try:
                raw_text = _extract_text_from_docx(file_path)
            except Exception:
                raw_text = ''
            if not raw_text.strip():
                raw_text = _extract_text_from_doc_ole(file_path)
        if not raw_text.strip():
            raise RuntimeError(f'未能从文件中提取到文字: {file_path}')
        result = structure_text(raw_text, DOC_PROMPT, api_key)
    else:
        raise ValueError(f'不支持的日程文件格式: {ext}（支持 jpg/png/pptx/ppt/docx/doc/xlsx）')

    return _parse_agenda_json(result)
