"""Extract background images from a PPT template using python-pptx.

No COM dependency — reads the PPTX file directly and searches for
background picture fills or large slide-covering images on key slides.
"""

import io
import os
import tempfile


def _extract_slide_images_ole(ppt_path: str, output_dir: str) -> tuple:
    """Fallback: extract the two largest images from a binary .ppt (OLE)
    file as template backgrounds."""
    import olefile
    import struct
    from PIL import Image

    paths = []

    try:
        ole = olefile.OleFileIO(ppt_path)
    except Exception:
        return None, None

    try:
        # Search for images in the Pictures stream
        if not ole.exists('Pictures'):
            return None, None
        data = ole.openstream('Pictures').read()

        sigs = [
            (b'\xff\xd8\xff', b'\xff\xd9', 'jpg'),
            (b'\x89PNG\r\n\x1a\n', b'IEND\xae\x42\x60\x82', 'png'),
        ]

        images = []
        for sig_bytes, end_marker, ext in sigs:
            start = 0
            while True:
                idx = data.find(sig_bytes, start)
                if idx == -1:
                    break
                end = data.find(end_marker, idx + 2)
                if end == -1:
                    end = len(data)
                else:
                    end += len(end_marker)
                img_data = data[idx:end]
                if len(img_data) > 10240:  # skip tiny images (<10KB)
                    try:
                        img = Image.open(io.BytesIO(img_data))
                        w, h = img.size
                        images.append((img_data, w * h, ext))
                    except Exception:
                        pass
                start = idx + 1

        # Pick the two largest images (likely full-slide backgrounds)
        images.sort(key=lambda x: x[1], reverse=True)
        if len(images) >= 2:
            home_path = os.path.join(output_dir, 'template_home.png')
            content_path = os.path.join(output_dir, 'template_content.png')
            _save_as_png_ole(images[0][0], home_path)
            _save_as_png_ole(images[1][0], content_path)
            return home_path, content_path
        return None, None
    finally:
        ole.close()


def _save_as_png_ole(blob, path):
    """Write image blob to *path* as a PNG."""
    from PIL import Image
    img = Image.open(io.BytesIO(blob))
    if img.mode in ('RGBA', 'P', 'LA', 'PA'):
        img = img.convert('RGB')
    img.save(path, 'PNG')


def extract_slide_images(pptx_path: str, output_dir: str = None) -> tuple:
    """Extract the home slide (slide 1) and a representative content slide
    (slide 4 if available, otherwise slide 2) as PNG background images.

    Supports .pptx (python-pptx) and .ppt (OLE fallback).

    Returns ``(home_img_path, content_img_path)``.
    """
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from PIL import Image

    if output_dir is None:
        output_dir = tempfile.mkdtemp(prefix='ppt_template_')

    # Try python-pptx first; binary .ppt will fail here
    try:
        prs = Presentation(pptx_path)
    except Exception:
        # Fall back to OLE extraction for binary .ppt
        home, content = _extract_slide_images_ole(pptx_path, output_dir)
        if home and content:
            return home, content
        raise RuntimeError(
            'PPT模板无法提取幻灯片背景图片（请确认模板包含图片背景）。'
            '\n.ppt 格式模板建议转换为 .pptx 格式后使用。'
        )

    total = len(prs.slides)

    def _image_from_blipfill(blipfill_el, part):
        """Return ``(blob, content_type)`` from an ``a:blipFill`` element
        using *part*'s relationships to resolve the ``r:embed`` id."""
        nsmap = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        }
        blips = blipfill_el.findall('.//a:blip', nsmap)
        for blip in blips:
            embed = blip.get(qn('r:embed'))
            if not embed or not hasattr(part, 'rels') or embed not in part.rels:
                continue
            rel = part.rels[embed]
            # rel is a _Relationship object; the actual image part is .target_part
            img_part = getattr(rel, 'target_part', rel)
            blob = getattr(img_part, 'blob', None) or getattr(img_part, '_blob', None)
            if blob:
                ct = getattr(img_part, 'content_type', 'image/png')
                return blob, ct
        return None, None

    def _bg_image_from_part(part, bg_xml):
        """Walk ``a:blipFill`` elements inside a background XML node."""
        nsmap = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }
        for bf in bg_xml.findall('.//a:blipFill', nsmap):
            blob, ct = _image_from_blipfill(bf, part)
            if blob:
                return blob, ct
        return None, None

    def _shape_has_picture_fill(shape, slide_part):
        """Check if a shape has a picture fill (blipFill) and return (blob, content_type)."""
        try:
            shape_xml = shape._element
            nsmap = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            }
            blipfills = shape_xml.findall('.//a:blipFill', nsmap)
            if not blipfills:
                return None, None
            for bf in blipfills:
                blob, ct = _image_from_blipfill(bf, slide_part)
                if blob:
                    return blob, ct
        except Exception:
            pass
        return None, None

    def _best_picture_shape(slide):
        """Return the image blob of the largest picture shape (or shape with
        picture fill) on the slide, or None."""
        best_blob = None
        best_area = 0

        for shape in slide.shapes:
            # Check direct picture shapes (shape_type == 13)
            if shape.shape_type == 13:
                area = shape.width * shape.height
                if area > best_area:
                    best_area = area
                    best_blob = (shape.image.blob, shape.image.content_type)
            else:
                # Check shapes with picture fill (common in templates: a
                # rectangle covering the whole slide with a picture fill)
                blob, ct = _shape_has_picture_fill(shape, slide.part)
                if blob:
                    area = shape.width * shape.height
                    if area > best_area:
                        best_area = area
                        best_blob = (blob, ct)

        return best_blob if best_blob else (None, None)

    def _extract_slide_bg(slide):
        """Extract the background image for a single slide.

        Priority:
        1. Slide-level background picture fill
        2. Slide-layout background picture fill
        3. Slide-master background picture fill
        4. Largest picture shape on the slide itself
        """
        # 1) Slide background
        try:
            bg_el = slide.background._element
            blob, ct = _bg_image_from_part(slide.part, bg_el)
            if blob:
                return blob, ct
        except Exception:
            pass

        # 2) Slide-layout background
        try:
            layout = slide.slide_layout
            layout_part = layout.part if hasattr(layout, 'part') else getattr(layout, '_part', None)
            if layout_part:
                layout_bg_el = layout.background._element
                blob, ct = _bg_image_from_part(layout_part, layout_bg_el)
                if blob:
                    return blob, ct
        except Exception:
            pass

        # 3) Slide-master background
        try:
            layout = slide.slide_layout
            master = layout.slide_master if hasattr(layout, 'slide_master') else None
            if master:
                master_part = master.part if hasattr(master, 'part') else getattr(master, '_part', None)
                if master_part:
                    master_bg_el = master.background._element
                    blob, ct = _bg_image_from_part(master_part, master_bg_el)
                    if blob:
                        return blob, ct
        except Exception:
            pass

        # 4) Largest picture shape
        return _best_picture_shape(slide)

    def _save_as_png(blob, path):
        """Write image blob to *path* as a PNG, converting to RGB if needed."""
        img = Image.open(io.BytesIO(blob))
        if img.mode in ('RGBA', 'P', 'LA', 'PA'):
            img = img.convert('RGB')
        img.save(path, 'PNG')

    # Main extraction
    home_path = os.path.join(output_dir, 'template_home.png')
    blob, _ct = _extract_slide_bg(prs.slides[0])
    if blob:
        _save_as_png(blob, home_path)

    # Content slide: slide 4 if available, otherwise slide 2 (0-indexed)
    content_idx = 3 if total >= 4 else min(1, total - 1)
    content_path = os.path.join(output_dir, 'template_content.png')
    blob2, _ct2 = _extract_slide_bg(prs.slides[content_idx])
    if blob2:
        _save_as_png(blob2, content_path)

    if not os.path.exists(home_path) or not os.path.exists(content_path):
        raise RuntimeError(
            'PPT模板无法提取幻灯片背景图片（幻灯片中未找到背景图或大图），'
            '请确认模板包含图片背景。'
        )

    return home_path, content_path
