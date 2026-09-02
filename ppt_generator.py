"""PPT generator with fixed layout. Two size configs: 16:9 and ultra-wide."""
from pptx import Presentation
from pptx.util import Inches, Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from extractors.speaker_extractor import Speaker, detect_face_strict
import os
import re

TEXT_COLOR = RGBColor(0x1F, 0x4E, 0x79)

# ---------------------------------------------------------------------------
# Layout configs — each size has its own set of positions and font sizes
# ---------------------------------------------------------------------------

CFG_16_9 = {
    'width': Inches(13.333),
    'height': Inches(7.5),
    'divider': {
        # 嘉宾姓名第一行约在页面中线（页面高 7.5"，中线 3.75"；
        # 文本框上边距 0.05"，speaker_top 3.7" → 文字起于 ≈3.75"）
        'title_top': Inches(2.5), 'title_h': Inches(0.8),
        'title_sz': Pt(40), 'title_color': TEXT_COLOR,
        'speaker_top': Inches(3.7), 'speaker_h': Inches(0.6),
        'speaker_sz': Pt(32), 'speaker_color': TEXT_COLOR,
    },
    'bio_photo': {
        'role_top': Inches(0.8), 'role_h': Inches(0.7), 'role_sz': Pt(36),
        # 照片上边沿与右侧履历文本框上沿对齐（bio_top = 1.5"）；
        # 下方姓名/机构说明文字随照片一起下移（name_top 保持照片下沿 + 0.14" 间距）
        'photo_left': Inches(0.8), 'photo_top': Inches(1.5),
        'photo_w': Inches(2.6), 'photo_h': Inches(3.8),
        'name_top': Inches(5.4), 'name_sz': Pt(18),
        'inst_top': Inches(5.8), 'inst_sz': Pt(14),
        'bio_left': Inches(4.0), 'bio_top': Inches(1.5),
        'bio_w': Inches(8.5), 'bio_h': Inches(4.85), 'bio_sz': Pt(17),
    },
    'bio_nophoto': {
        'role_top': Inches(0.8), 'role_h': Inches(0.7), 'role_sz': Pt(36),
        'bio_left': Inches(1.5), 'bio_top': Inches(1.8),
        'bio_w': Inches(10.3), 'bio_h': Inches(4.65), 'bio_sz': Pt(17),
    },
    'topic': {
        'title_top': Inches(3.0), 'title_h': Inches(0.8),
        'title_sz': Pt(40), 'title_color': TEXT_COLOR,
        'speaker_top': Inches(3.9), 'speaker_h': Inches(1.6),
        'speaker_sz': Pt(36), 'speaker_color': TEXT_COLOR,
    },
    'countdown': {
        'cn_top': Inches(2.5), 'cn_sz': Pt(20),
        'en_top': Inches(4.5), 'en_sz': Pt(14),
    },
}

CFG_ULTRAWIDE = {
    'width': Inches(7.874),
    'height': Inches(1.575),
    'divider': {
        'title_top': Inches(0.45), 'title_h': Inches(0.6),
        'title_sz': Pt(32), 'title_color': TEXT_COLOR,
        'speaker_top': Inches(0.6), 'speaker_h': Inches(0.45),
        'speaker_sz': Pt(22), 'speaker_color': TEXT_COLOR,
    },
    'bio_photo': {
        'role_top': Inches(0.05), 'role_h': Inches(0.3), 'role_sz': Pt(18),
        'photo_left': Inches(0.1), 'photo_top': Inches(0.35),
        'photo_w': Inches(1.5), 'photo_h': Inches(1.15),
        'name_top': Inches(0.55), 'name_sz': Pt(10),
        'inst_top': Inches(0.85), 'inst_sz': Pt(8),
        'bio_left': Inches(1.8), 'bio_top': Inches(0.35),
        'bio_w': Inches(5.8), 'bio_h': Inches(0.96), 'bio_sz': Pt(9),
    },
    'bio_nophoto': {
        'role_top': Inches(0.05), 'role_h': Inches(0.3), 'role_sz': Pt(18),
        'bio_left': Inches(0.3), 'bio_top': Inches(0.4),
        'bio_w': Inches(7.2), 'bio_h': Inches(0.89), 'bio_sz': Pt(10),
    },
    'topic': {
        'title_top': Inches(0.3), 'title_h': Inches(0.55),
        'title_sz': Pt(28), 'title_color': TEXT_COLOR,
        'speaker_top': Inches(0.7), 'speaker_h': Inches(0.55),
        'speaker_sz': Pt(20), 'speaker_color': TEXT_COLOR,
    },
    'countdown': {
        'cn_top': Inches(0.2), 'cn_sz': Pt(10),
        'en_top': Inches(0.65), 'en_sz': Pt(7),
    },
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _split_names(name_str):
    """Split a name string into individual names. Handles multi-person entries
    and keeps English names (e.g. "Theo M. de Reijke") as one person."""
    if not name_str or not name_str.strip():
        return []
    # Split by common delimiters: Chinese comma, English comma, semicolon, slash, etc.
    parts = re.split(r'[,，、/；;]+|\s+和\s+|\s+与\s+', name_str)
    # Also split "王艳波 教授 王凯臣 教授" by whitespace, dropping title tokens
    titles = {'教授', '主任', '主任医师', '副主任医师', '主治医师', '医师',
              '院长', '副院长', '博士', '院士', '研究员', '副研究员', '主席',
              'Prof', 'Prof.', 'Dr', 'Dr.', 'MD', 'PhD', 'Mr', 'Mr.', 'Ms', 'Ms.'}
    names = []
    for p in parts:
        buf = []  # pending English-name tokens, joined back into one name
        for token in re.split(r'\s+', p.strip()):
            if not token:
                continue
            if token in titles:
                if buf:
                    names.append(' '.join(buf))
                    buf = []
                continue
            if re.fullmatch(r'[一-鿿]{1,4}', token):
                # Chinese name token — starts a new person
                if buf:
                    names.append(' '.join(buf))
                    buf = []
                names.append(token)
            else:
                # English name fragment (e.g. "Theo", "M.", "de", "Reijke")
                buf.append(token)
        if buf:
            names.append(' '.join(buf))
    return names


def _spread_name(name):
    """Insert a non-breaking space in 2-char Chinese names for alignment:
    陈晨 → 陈 晨（U+00A0，防止姓名在换行处被拆开）"""
    if len(name) == 2 and re.fullmatch(r'[一-鿿]{2}', name):
        return f'{name[0]}\u00A0{name[1]}'
    return name


def _match_speaker(speaker_name, speakers):
    """Fuzzy-match an agenda speaker_name to a key in speakers dict."""
    if not speaker_name:
        return None
    # Exact match first
    if speaker_name in speakers:
        return speakers[speaker_name]
    # Strip common titles and middle initials (e.g., "M.", "J.")
    clean = re.sub(r'(Prof\.?|Dr\.?|Mr\.?|Ms\.?|PhD|MD)\s*', '', speaker_name,
                   flags=re.IGNORECASE).strip()
    clean = re.sub(r'\b[A-Z]\.\s*', '', clean).strip()
    if clean in speakers:
        return speakers[clean]
    # Substring matching: agenda name contains speaker key, or vice versa
    for key in speakers:
        if len(key) >= 2 and (key in speaker_name or speaker_name in key
                              or key in clean or clean in key):
            return speakers[key]
    return None

def _force_font_theme(prs):
    """Override theme font scheme to Microsoft YaHei."""
    FONT = 'Microsoft YaHei'
    from lxml import etree

    for master in prs.slide_masters:
        for rel in master.part.rels.values():
            if 'theme' not in rel.reltype:
                continue
            root = etree.fromstring(rel.target_part.blob)
            for elem in root.iter():
                tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                if tag in ('ea', 'latin', 'cs') and elem.get('typeface'):
                    elem.set('typeface', FONT)
            rel.target_part._blob = etree.tostring(root)
            break
        break


def _add_bg(slide, img_path, w, h):
    """Add a full-slide background image."""
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, 0, 0, w, h)


def _set_font_ea(paragraph, font_name):
    """Force East Asian font on a paragraph's first run."""
    if not paragraph.runs:
        return
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    rPr = paragraph.runs[0]._r.get_or_add_rPr()
    for tag in ('ea', 'latin', 'cs'):
        for old in rPr.findall(f'{{{A_NS}}}{tag}'):
            rPr.remove(old)
    ea = rPr.makeelement(f'{{{A_NS}}}ea', {'typeface': font_name})
    rPr.insert(0, ea)
    latin = rPr.makeelement(f'{{{A_NS}}}latin', {'typeface': font_name})
    rPr.insert(0, latin)


def _apply_src_rect(pic, l, t, r, b):
    """Apply source rectangle crop to a picture shape (values in 1/1000 percent)."""
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    srcRect = pic._element.makeelement(f'{{{A_NS}}}srcRect', {})
    blipFill = pic._element.find(f'{{{P_NS}}}blipFill')
    if blipFill is not None:
        blipFill.insert(1, srcRect)
    if l:
        srcRect.set('l', str(int(l)))
    if t:
        srcRect.set('t', str(int(t)))
    if r:
        srcRect.set('r', str(int(r)))
    if b:
        srcRect.set('b', str(int(b)))


def _add_textbox(slide, left, top, width, height, text, font_size, color=None,
                 bold=False, align=PP_ALIGN.CENTER, font_name='Microsoft YaHei'):
    if color is None:
        color = TEXT_COLOR
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    _set_font_ea(p, font_name)
    return txBox


def _add_multiline_textbox(slide, left, top, width, height, text_lines,
                           font_size, color=None, align=PP_ALIGN.LEFT,
                           font_name='Microsoft YaHei', line_spacing=1.2,
                           bullet=False, bold=False):
    if color is None:
        color = TEXT_COLOR
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(line_spacing)
        _set_font_ea(p, font_name)
        if bullet:
            pPr = p._p.get_or_add_pPr()
            buChar = pPr.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}buChar',
                                     {'char': '•'})
            for child in list(pPr):
                if child.tag.endswith('}buChar') or child.tag.endswith('}buNone'):
                    pPr.remove(child)
            pPr.append(buChar)
            # Add indent: ~1 character space between bullet and text
            pPr.set('indent', '355600')  # ~0.5cm / ~1 char in EMU
            pPr.set('marL', '355600')     # left margin to align with indent
    return txBox


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _make_cover(prs, home_bg, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(slide, home_bg, cfg['width'], cfg['height'])
    return slide


def _make_countdown(prs, home_bg, cfg, cn_text, en_text, bilingual):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, home_bg, cfg['width'], cfg['height'])
    c = cfg['countdown']
    w = cfg['width']
    _add_textbox(slide, Inches(0.5), c['cn_top'], w - Inches(1), Inches(0.6),
                 cn_text, c['cn_sz'])
    if bilingual:
        _add_textbox(slide, Inches(0.5), c['en_top'], w - Inches(1), Inches(0.4),
                     en_text, c['en_sz'])
    return slide


def _make_divider(prs, home_bg, cfg, title_cn, title_en, bilingual, speaker_line=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, home_bg, cfg['width'], cfg['height'])
    d = cfg['divider']
    w = cfg['width']

    if bilingual and title_en:
        full_title = f'{title_cn}  /  {title_en}'
    else:
        full_title = title_cn
    _add_textbox(slide, Inches(0.5), d['title_top'], w - Inches(1), d['title_h'],
                 full_title, d['title_sz'], color=d['title_color'], bold=True)

    if speaker_line:
        # If title is long (>22 chars), shift speaker down by one title line height
        speaker_top = d['speaker_top']
        if len(full_title) > 22:
            speaker_top += d['title_sz']
        line_count = speaker_line.count('\n') + 1
        speaker_h = d['speaker_h'] * line_count
        # 嘉宾介绍行：文本框左右各缩短 2 字符（1 字符宽 ≈ 1 个字号），
        # 盒子相对页面中线对称，居中显示
        left = Inches(0.5) + 2 * d['speaker_sz']
        _add_textbox(slide, left, speaker_top, w - 2 * left, speaker_h,
                     speaker_line, d['speaker_sz'], color=d['speaker_color'],
                     bold=True, align=PP_ALIGN.CENTER)
    return slide


ROLE_LABELS_EN = {
    '大会主席': 'Chairman',
    '分会场主席': 'Session Chair',
    '总结嘉宾': 'Closing Remarks',
    '主讲嘉宾': 'Speaker',
    '主持嘉宾': 'Moderator',
    '嘉宾': 'Guest',
}


def _make_bio(prs, content_bg, cfg, speaker, role_label='', en=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, content_bg, cfg['width'], cfg['height'])
    w = cfg['width']

    has_photo = bool(speaker.photo_path and os.path.exists(speaker.photo_path))
    bp = cfg['bio_photo'] if has_photo else cfg['bio_nophoto']

    if en:
        bio_text = speaker.bio_en
        name_display = f'Prof. {speaker.name}'
        if role_label:
            role_label = ROLE_LABELS_EN.get(role_label, role_label)
    else:
        bio_text = speaker.bio
        name_display = speaker.name
        if speaker.title:
            name_display = f'{speaker.name} {speaker.title}'

    if role_label:
        _add_textbox(slide, bp['bio_left'], bp['role_top'], bp['bio_w'], bp['role_h'],
                     role_label, Pt(36), bold=True, align=PP_ALIGN.CENTER)

    if has_photo:
        TARGET_RATIO = 0.68
        photo_width = Cm(6.5)
        try:
            from PIL import Image
            img = Image.open(speaker.photo_path)
            pw, ph = img.size
            w_h_ratio = pw / ph if ph > 0 else 1

            # Detect face for smart centering（多级联严格投票，避免误检/合照误定位）
            face_cx = pw / 2
            face_cy = ph * 0.4
            face_found = False
            try:
                import cv2, numpy as np
                pil_img = img.convert('RGB')
                cv_img = cv2.cvtColor(np.array(pil_img), cv2.COLOR_RGB2BGR)
                face = detect_face_strict(cv_img)
                if face is not None:
                    x, y, fw, fh = face
                    face_cx = x + fw / 2
                    face_cy = y + fh / 2
                    face_found = True
            except Exception:
                pass

            # Calculate visible region centered on face, constrained to image bounds
            # 全身照 → 上半身显示裁剪；区域放不下时退回原比例适配（保证不出界）
            upper_ok = False
            if face_found and fh < 0.4 * ph:
                # 可见高度 4.5×人脸高（脸约占 22%，上留头发、下至胸肩）
                target_h = int(4.5 * fh)
                target_w = int(target_h * TARGET_RATIO)
                if target_w <= pw and target_h <= ph:
                    anchor_x = face_cx - target_w / 2
                    anchor_y = y - 1.2 * fh
                    upper_ok = True
            if not upper_ok:
                if w_h_ratio < TARGET_RATIO:
                    # Too tall: keep full width, crop height
                    target_w = pw
                    target_h = int(pw / TARGET_RATIO)
                else:
                    # Too wide: keep full height, crop width
                    target_w = int(ph * TARGET_RATIO)
                    target_h = ph
                anchor_x = face_cx - target_w / 2
                anchor_y = face_cy - target_h * 0.4

            vis_l = max(0, min(int(anchor_x), pw - target_w))
            vis_t = max(0, min(int(anchor_y), ph - target_h))
            vis_r = vis_l + target_w
            vis_b = vis_t + target_h

            crop_l = int(vis_l / pw * 100000)
            crop_t = int(vis_t / ph * 100000)
            crop_r = int((pw - vis_r) / pw * 100000)
            crop_b = int((ph - vis_b) / ph * 100000)

            photo_height = int(photo_width / TARGET_RATIO)
            pic = slide.shapes.add_picture(speaker.photo_path,
                                           bp['photo_left'], bp['photo_top'],
                                           photo_width, photo_height)

            if crop_l or crop_r or crop_t or crop_b:
                _apply_src_rect(pic, crop_l, crop_t, crop_r, crop_b)

            # Add shadow effect: transparency 35%, blur 23pt, distance 11pt
            spPr = pic._element.find(
                '{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
            if spPr is None:
                spPr = pic._element.makeelement(
                    '{http://schemas.openxmlformats.org/presentationml/2006/main}spPr', {})
                pic._element.insert(0, spPr)
            effectLst = spPr.makeelement(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst', {})
            outerShdw = effectLst.makeelement(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw', {
                    'blurRad': str(int(23 * 12700)),
                    'dist': str(int(11 * 12700)),
                    'dir': '8100000',
                })
            srgbClr = outerShdw.makeelement(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr',
                {'val': '000000'})
            alpha = srgbClr.makeelement(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha',
                {'val': '35000'})
            srgbClr.append(alpha)
            outerShdw.append(srgbClr)
            effectLst.append(outerShdw)
            spPr.append(effectLst)
        except Exception:
            pass

        # Name + title on line 1, institution on line 2
        info_lines = [name_display]
        if speaker.institution:
            info_lines.append(speaker.institution)
        if info_lines:
            _add_multiline_textbox(slide, bp['photo_left'] - Cm(1.2), bp['name_top'],
                                   Cm(8.9), Inches(0.7),
                                   info_lines, Pt(18),
                                   align=PP_ALIGN.CENTER, bold=True)

    # Bio text (split into lines)
    raw_lines = [l.strip() for l in bio_text.replace('\r', '').split('\n') if l.strip()]
    if not raw_lines:
        # Speaker bio is missing — show fallback text
        if en:
            fallback = f'{speaker.name} (No information available)'
        else:
            fallback = f'{speaker.name}（信息缺失）'
        _add_textbox(slide, bp['bio_left'], bp['bio_top'],
                     bp['bio_w'], Inches(1.2),
                     fallback, Pt(20),
                     color=TEXT_COLOR, align=PP_ALIGN.CENTER)
        return slide

    # Split very long lines (>=100 chars) at Chinese delimiters 、or ，
    bio_lines = []
    for line in raw_lines:
        if len(line) >= 100:
            # Split at 、and ，, discarding the delimiters
            parts = []
            cur = ''
            for ch in line:
                if ch in ('、', '，'):
                    if cur.strip():
                        parts.append(cur.strip())
                    cur = ''
                else:
                    cur += ch
            if cur.strip():
                parts.append(cur.strip())
            if parts:
                bio_lines.extend(parts)
            else:
                bio_lines.append(line)
        else:
            bio_lines.append(line)

    # Effective visual lines: >=28 chars per line wraps to 2
    N_visual = sum(2 if len(line) >= 28 else 1 for line in bio_lines)

    # Font size: shrink for very long bios
    if N_visual >= 20:
        bio_sz = Pt(14)
        FS_pt = 14
    elif N_visual >= 15:
        bio_sz = Pt(16)
        FS_pt = 16
    else:
        bio_sz = bp['bio_sz']
        FS_pt = bio_sz / 12700

    # # Fixed line spacing rules based on effective line count
    # # (multiplier × default line height ≈ 1.2×FS; extra = (multiplier−1) × 1.2×FS)
    # if N_visual <= 8:
    #     bio_spacing = FS_pt * 1.2     # 2× spacing
    # elif N_visual <= 12:
    #     bio_spacing = FS_pt * 0.6     # 1.5× spacing
    # else:
    #     bio_spacing = 0.12               # 1.1× spacing (single)
    # 基础偏移量定为 1.2，让 8 行时大约是 2倍，12行时大约是 1.5倍
    bio_spacing = FS_pt * max(0, 2.4 - 0.15 * N_visual) # 用max确保不会变成负数

    _add_multiline_textbox(slide, bp['bio_left'], bp['bio_top'],
                           bp['bio_w'], bp['bio_h'],
                           bio_lines, bio_sz, bullet=True,
                           line_spacing=bio_spacing)
    return slide


def _make_topic(prs, content_bg, cfg, title_cn, title_en, bilingual,
                speaker_name, institution):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, content_bg, cfg['width'], cfg['height'])
    t = cfg['topic']
    w = cfg['width']

    if bilingual and title_en:
        title_line = f'{title_cn}  |  {title_en}'
    else:
        title_line = title_cn
    _add_textbox(slide, Inches(0.3), t['title_top'], w - Inches(0.6), t['title_h'],
                 title_line, t['title_sz'], color=t['title_color'], bold=True)

    # Insert space in 2-char Chinese names for visual balance
    def _pad_name(n):
        return f'{n[0]}  {n[1]}' if len(n) == 2 and '一' <= n[0] <= '鿿' else n

    speaker_line = _pad_name(speaker_name)
    if institution:
        speaker_line = f'{speaker_line}  |  {institution}'
    _add_textbox(slide, Inches(0.3), t['speaker_top'], w - Inches(0.6), t['speaker_h'],
                 speaker_line, t['speaker_sz'], color=t['speaker_color'], bold=True)
    return slide


def _make_transition(prs, content_bg, cfg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, content_bg, cfg['width'], cfg['height'])
    return slide


# ---------------------------------------------------------------------------
# Main entry
# ---------------------------------------------------------------------------

def generate_ppt(agenda_items, speakers, home_bg, content_bg,
                 slide_size, lang, output_path):
    """
    Generate the full host-script PPT.

    agenda_items: list of AgendaItem
    speakers: dict of {name: Speaker}
    home_bg: path to home/divider background image
    content_bg: path to content page background image
    slide_size: '16:9' or 'ultrawide'
    lang: 'bilingual' or 'chinese'
    output_path: where to save the .pptx
    """
    bilingual = (lang == 'bilingual')
    cfg = CFG_16_9 if slide_size == '16:9' else CFG_ULTRAWIDE

    prs = Presentation()
    prs.slide_width = cfg['width']
    prs.slide_height = cfg['height']

    # Force Microsoft YaHei at theme/master level
    _force_font_theme(prs)

    # 1. Cover
    _make_cover(prs, home_bg, cfg)

    # 2. Opening countdown slides (use content page background)
    _make_countdown(prs, content_bg, cfg,
                    '尊敬的各位来宾，本次会议还有【5分钟】开始，请您尽快就坐，'
                    '并将手机等通讯设备关闭或置于静音状态，谢谢您的合作。',
                    'Dear guests, the meeting will begin in [5 minutes]. '
                    'Please take your seats promptly and turn off or silence '
                    'your mobile phones and other communication devices. Thank you.',
                    bilingual)
    _make_countdown(prs, content_bg, cfg,
                    '尊敬的各位来宾，本次会议即将开始，请您尽快就坐，'
                    '并将手机等通讯设备关闭或置于静音状态，谢谢您的合作。',
                    'Dear guests, the meeting is about to begin. '
                    'Please take your seats promptly and turn off or silence '
                    'your mobile phones and other communication devices. Thank you.',
                    bilingual)

    # 3. Per-agenda-item slide groups
    for item in agenda_items:
        # Divider: format speaker/host lines with role labels
        role_map_divider = {
            'opening': '大会主席',
            'closing': '总结嘉宾',
            'speech': '主讲嘉宾',
            'panel': '主讲嘉宾',
        }
        speaker_role = role_map_divider.get(item.item_type, '主讲嘉宾')

        def _format_name(name):
            sp = _match_speaker(name, speakers)
            display = sp.name if sp else name
            if re.search(r'[A-Za-z]', display):
                # English name — prefix with "Prof. "（不换行空格，整段姓名不拆行）
                return 'Prof.\u00A0' + display.replace(' ', '\u00A0')
            if sp and sp.title:
                return f'{_spread_name(display)}\u00A0{sp.title}'
            return f'{_spread_name(display)}\u00A0教授'

        def _format_persons(name_str, role_label):
            if not name_str:
                return ''
            parts = [_format_name(name) for name in _split_names(name_str)]
            return f'{role_label}：{"、".join(parts)}' if parts else ''

        if item.item_type in ('opening', 'closing'):
            # Opening/closing divider: merge speakers + hosts, dedupe by
            # name, show "姓名 职称" only — no role labels
            seen = set()
            parts = []
            for name_str in (item.speaker_name, item.host):
                for name in _split_names(name_str):
                    if name in seen:
                        continue
                    seen.add(name)
                    parts.append(_format_name(name))
            speaker_line = '、'.join(parts)
        else:
            lines = []
            sp_line = _format_persons(item.speaker_name, speaker_role)
            if sp_line:
                lines.append(sp_line)
            host_line = _format_persons(item.host, '主持嘉宾')
            if host_line:
                lines.append(host_line)
            speaker_line = '\n'.join(lines)
        _make_divider(prs, content_bg, cfg,
                      item.session_title_cn, item.session_title_en,
                      bilingual, speaker_line)

        # Section headers only need a divider slide, skip bio and topic
        if item.item_type == 'section':
            continue

        # Collect all persons: speakers + hosts, each gets a bio page
        role_map = {
            'opening': '大会主席',
            'closing': '总结嘉宾',
            'speech': '主讲嘉宾',
            'panel': '主讲嘉宾',
        }
        role = role_map.get(item.item_type, '嘉宾')
        # 分会场主席致辞 → 分会场主席
        title_text = item.session_title_cn + item.session_title_en
        if '分会' in title_text and item.item_type in ('opening', 'speech'):
            role = '分会场主席'
        host_role = '主持嘉宾'

        persons = []  # [(name, role_label)]
        for name in _split_names(item.speaker_name):
            persons.append((name, role))
        for name in _split_names(item.host):
            persons.append((name, host_role))

        # Opening/closing: each person gets only one bio page
        # (speaker and host lists often overlap, e.g. 主席致辞)
        if item.item_type in ('opening', 'closing'):
            seen = set()
            deduped = []
            for name, label in persons:
                if name in seen:
                    continue
                seen.add(name)
                deduped.append((name, label))
            persons = deduped

        for person_name, label in persons:
            sp = _match_speaker(person_name, speakers)
            if sp:
                _make_bio(prs, content_bg, cfg, sp, label)
                # English-named speakers get a second bio page in English
                if re.search(r'[A-Za-z]', sp.name) and sp.bio_en:
                    _make_bio(prs, content_bg, cfg, sp, label, en=True)
            else:
                _make_bio(prs, content_bg, cfg,
                          Speaker(name=person_name, bio=item.institution), label)

    # 4. Closing — blank page with home background only
    _make_cover(prs, home_bg, cfg)

    prs.save(output_path)
    return output_path


# ---------------------------------------------------------------------------
# CLI entry — called via subprocess from main_window.py
# ---------------------------------------------------------------------------

if __name__ == '__main__':
    import sys
    import json as _json

    if len(sys.argv) < 2:
        print('Usage: python ppt_generator.py <params_json_path>', file=sys.stderr)
        sys.exit(1)

    params_path = sys.argv[1]
    with open(params_path, 'r', encoding='utf-8') as f:
        params = _json.load(f)

    # Reconstruct AgendaItem objects
    agenda_items = []
    for d in params['agenda_items']:
        agenda_items.append(__import__('extractors.agenda_extractor',
                                       fromlist=['AgendaItem']).AgendaItem(
            order=d['order'],
            time_slot=d['time_slot'],
            session_title_cn=d['session_title_cn'],
            session_title_en=d['session_title_en'],
            speaker_name=d['speaker_name'],
            host=d['host'],
            institution=d['institution'],
            item_type=d['item_type'],
        ))

    # Reconstruct Speaker objects
    speakers = {}
    for name, d in params['speakers'].items():
        speakers[name] = Speaker(
            name=d['name'],
            photo_path=d['photo_path'],
            bio=d['bio'],
            bio_en=d.get('bio_en', ''),
            institution=d['institution'],
            title=d['title'],
        )

    generate_ppt(
        agenda_items=agenda_items,
        speakers=speakers,
        home_bg=params['home_bg'],
        content_bg=params['content_bg'],
        slide_size=params['slide_size'],
        lang=params['lang'],
        output_path=params['output_path'],
    )
